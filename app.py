import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
import asyncio

from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai

from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

# For Streamlit Deployement Load environment variables
genai_api_key = st.secrets["GOOGLE_API_KEY"]
genai.configure(api_key=genai_api_key)

# Function to extract text from PDF files
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Function to extract text from PPTX files
def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Function to create and save vector store
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Function to create conversational chain
async def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, 
    if the answer is not in the provided context just say,"answer is not available in the context", 
    don't provide the wrong answer
    1.	Income Tax Rules from :

What is Income Tax?
Income tax is a type of direct tax paid to the central government of India. It is the primary source of revenue for the government. It plays an important role in financing various developmental and welfare activities. It is governed by the Income Tax Act of 1961, and is administered by the Central Board of Direct Taxes (CBDT) under the Ministry of Finance, Government of India. there are various types of income tax like

Personal Tax: It applies to the income earned by the individual and HUFs. It is calculated by the individual’s total taxable income sourced from salary, house property, capital gain, etc.

Corporate Income Tax: It applies to the companies registered under the Companies Act, 2013.

Capital Gains Tax

Securities Transaction Tax (STT)

Dividend Distribution Tax (DDT)

Minimum Alternate Tax (MAT)

Tax Deducted at Source

Interim Budget 2024 Updates:
1. Maintenance of Existing Tax Rates: The interim budget for 2024 keeps the current tax rates unchanged for both direct and indirect taxes.

2. Tax Exemption for Lower Income Groups: Individuals earning up to Rs 7 lakh annually have no tax liability, providing relief for lower income groups.

3. Withdrawal of Tax Dispute Provisions: Finance Minister Nirmala Sitharaman withdraws tax dispute claims up to Rs 25,000 for disputes pertaining to the financial year 2009-10 and Rs 10,000 for disputes from financial years 2010-11 to 2014-15.

4. Enhanced Leave Encashment Limit: The limit for leave encashment for non-government employees has been raised significantly from Rs 3 lakh to Rs 25 lakh.

5. Reduction in TDS Rate on EPF Withdrawal: The Tax Deducted at Source (TDS) rate on Employee Provident Fund (EPF) withdrawals has been decreased from 30% to 20%, providing relief to EPF subscribers.

6. Standard Deduction for Salaried Employees: Salaried employees and pensioners can claim a standard deduction of Rs 50,000 under the new tax regime.

7. Reduction in Highest Surcharge: The highest surcharge for individuals earning more than Rs 5 crore has been reduced to 25% from the previous 37%, resulting in a decreased tax rate for this bracket.

8. Opt-Out Provision for New Tax Regime: While the new tax regime will be the default option, taxpayers have the option to opt out before the due date for filing income tax returns for the respective assessment year.




Budget 2023 Major Updates:
1. Tax Rebate for Lower Income Groups: Individuals earning up to Rs 7 lakh annually receive a tax rebate in the New Tax Regime, exempting them from paying taxes if their taxable income is below Rs 7 lakh.

 

The new tax slabs under the new tax regime will be:
Income Slabs

Tax Rates

up to Rs 3 lakh

Nil

Rs 3 lakh- Rs 6 lakh

5%

Rs 6 lakh-Rs 9 lakh

10%

Rs 9 lakh-Rs 12 lakh

15%

Rs 12 lakh- Rs 15 lakh

20%

Above Rs 15 lakh

30%

 

 

2. Standard Deduction for Salaried Employees and Pensioners: Salaried employees and pensioners can claim a standard deduction of Rs 50,000 under the new tax regime, reducing their taxable income.

3. Reduction in Highest Surcharge: The highest surcharge for individuals earning over Rs 5 crore has been lowered to 25% from 37%, resulting in a reduced tax rate of 39% for this group.

4. Default New Tax Regime: The new Income Tax (IT) regime becomes the default tax regime, providing a simplified tax structure for taxpayers. However, individuals have the option to opt out before the due date for filing IT returns for the respective assessment year.

5. Enhanced Leave Encashment Limit: Non-government employees benefit from an increased limit for leave encashment, raised from Rs 3 lakh to Rs 25 lakh, offering greater flexibility in utilizing accrued leave benefits.

6. Decrease in TDS Rate on EPF Withdrawal: The Tax Deducted at Source (TDS) rate on Employee Provident Fund (EPF) withdrawals has been reduced from 30% to 20%, easing the tax burden on EPF subscribers during withdrawal.




Who should pay Income Tax?
Income Tax applies to individuals, businesses and other entities as per their income or profit earned in a particular year.

Different tax rules apply to different types of taxpayers.

   Below are the categories of taxpayers:

Individuals
Hindu Undivided Family (HUF)
Firms
Companies
Association of Persons(AOP)
Body of Individuals (BOI)
Local Authority
Artificial Judicial Person
Here are some point about who should pay tax-

Individual
Individuals who are residents of India need to pay income tax on their total taxable income.

Non-resident individuals also need to pay taxes earned or received in India.

Income tax is applicable on sources like salaries, house property, capital gains, business or profession, and other incomes like interest, dividends, etc.

The tax rates vary as per age, tax slabs, etc.

Hindu Undivided Families (HUFs)
HUFs are also eligible to pay income tax on their total taxable income.

HUFs are taxed separately from the income of individual members of a family.

Companies, Partnership & LLPs
Companies, Partnerships and Limited Liability Partnerships (LLPs) registered under the Companies Act, 2013 need to pay income tax on the profit gained in the particular financial year.

The tax rate is based on the turnover, business structure, etc.

Other Entities
Trusts, associations, institutions, and other legal entities are also liable to pay income tax as per their income and source.

Types of Income – What are the 5 heads of income?
All individuals earning income in India, whether residents or non-residents, are liable to pay income tax. To simplify categorization, the Income Tax Department divides income into five primary heads:




Head of Income and Nature of Income Covered:
Head of Income	Nature of Income Covered
Income from Other Sources	Taxable income from sources such as interest earned on savings bank accounts, fixed deposits, and winnings from lotteries.
Income from House Property	Taxable income derived from renting out a property, including residential and commercial spaces.
Income from Capital Gains	Taxable surplus generated from the sale of capital assets like mutual funds, shares, and real estate properties.
Income from Business and Profession	Taxable profits earned by self-employed individuals, freelancers, contractors, and professionals like doctors, lawyers, and insurance agents.
Income from Salary	Taxable income earned from employment, including salaries and pensions.
Simplified Explanation:
Income from Other Sources: Tax on money earned from savings account interest, fixed deposits, or lottery winnings.
Income from House Property: Tax on money earned from renting out properties, whether residential or commercial.
Income from Capital Gains: Tax on profits made by selling investments like mutual funds, stocks, or real estate.
Income from Business and Profession: Tax on profits earned by self-employed individuals, freelancers, and professionals like doctors or lawyers.
Income from Salary: Tax on money earned through employment, including salaries and pensions.
Components of Income Tax
Here is the key component of income tax, that you should know

Gross Total Income (GTI): Gross Total Income is the total income earned by the individual business or other entities during a financial year before any deduction or exemption. It includes salaries, wages, profit, capital gains, house property, etc.

Deductions: It is expenses or investments that taxpayers can claim to reduce their income amount. There are various sections in the Income Tax Act like Section 80C, Section 80D, Section 80 G, etc. It lowers the tax liability and encourages saving and investments.

Taxable Income: It is the income remaining after deducting a deduction from the GTI. It is the amount on which income tax is calculated and paid to the government.

Tax Slabs and Rates: Income tax in India is levied at progressive tax rates, which means higher incomes attract high-income tax rates. It is revised every year by the government through the union budget.

Tax Credit: Tax Credits are incentives provided by the government. Common tax credits include deductions for taxes paid in advance (TDS), taxes paid on self-assessment, foreign tax credits for taxes paid in another country, etc.

Tax Liability: It is the total amount of tax needed to be paid by the taxpayer to the government. It is calculated by applying the applicable tax rates to the taxable income after considering deductions and tax credits.

What is the Income Tax Regime?
Income Tax Regime is the system and framework established by the government to levy taxes on the income earned by individuals, businesses, and other entities within the country. It is a tax slab by the government.

Latest New tax regime FY 2023-24
Income up to Rs 3 lakh: Nil

Income Rs 3 lakh to Rs 6 lakh: 5%

Income Rs 6 lakh to Rs 9 lakh: 10%

Income Rs 9 lakh to Rs 12 lakh: 15%

Income Rs 12 lakh to Rs 15 lakh: 20%

Income above Rs 15 lakh: 30%

 

Last Tax Regime FY 2022-23
Income up to Rs 2.5 lakh: Nil

Income Rs 2.5 lakh to Rs 5 lakh: 5%

Income Rs 5 lakh to Rs 7.5 lakh: 10%

Income Rs 7.5 lakh to Rs 10 lakh: 15%

Income Rs 10 lakh to Rs 12.5 lakh: 20%

Income Rs 12.5 lakh to Rs 15 lakh: 25%

Income above Rs 15 lakh: 30%

Taxpayers and Tax Slabs
– Taxpayers in India are divided into different categories like individuals, Hindu Undivided Families (HUFs), associations of persons (AOPs), and bodies of individuals (BOIs).
– Tax rates for firms and Indian companies are fixed, whereas individual taxpayers are taxed based on income slabs.
– Income is categorized into blocks called tax slabs, each with a different tax rate.
– As income increases, the tax rate also increases, following a progressive tax system.
– This means that higher income earners pay a higher percentage of their income in taxes.

What is the Existing/Old Income Tax Regime?
– The old income tax regime had three slab rates: 5%, 20%, and 30%.
– Taxpayers could opt for this regime and claim various deductions.
– Allowances like Leave Travel Concession (LTC), House Rent Allowance (HRA), and specific other allowances were deductible.
– Deductions for tax-saving investments under sections 80C to 80U were allowed.
– A standard deduction of Rs 50,000 was applicable.
– Taxpayers could also claim a deduction for interest paid on home loans.

Tax slab rates applicable for Individual taxpayers below 60 years for the Old tax regime are as below:

Income Range

Tax rate

Tax to be paid

Up to Rs 2,50,000

0

No tax

Rs 2.5 lakhs – Rs 5 lakhs

5%

5% of your taxable income

Rs 5 lakhs – Rs 10 lakhs

20%

Rs 12,500+20% on income above Rs 5 lakh

Above 10 lakhs

30%

Rs 1,12,500+30% on income above Rs 10 lakh




There are additional tax slabs for individuals aged 60 and above, as well as those above 80.
It’s worth noting that some people mistakenly assume that if they earn Rs. 12 lakh, they’ll pay 30% tax on the entire amount (Rs. 3,60,000). However, under the progressive tax system, someone earning Rs. 12 lakh would pay Rs. 1,12,500 + Rs. 60,000, totaling Rs. 1,72,500.

Income Tax Slabs Under New Tax Regime
In the 2020 budget, a new tax system was introduced with reduced tax rates and fewer deductions for Individuals and HUFs. As a result, many taxpayers didn’t choose the new system. However, to promote it in Budget 2023, the income tax slabs for FY 2023-24 (AY 2024-25) have been updated as mentioned below:-

New tax regime FY 2023-24
(After budget)

New tax regime FY 2022-23
(Before budget)

Income up to Rs 3 lakh

Nil

Up to Rs 2.5 lakh

Nil

Rs 3 lakh to Rs 6 lakh

5%

Rs 2.5 lakh to Rs 5 lakh

5%

Rs 6 lakh to Rs 9 lakh

10%

Rs 5 lakh to Rs 7.5 lakh

10%

Rs 9 lakh to Rs 12 lakh

15%

Rs 7.5 lakh to Rs 10 lakh

15%

Rs 12 lakh to Rs 15 lakh

20%

Rs 10 lakh to Rs 12.5 lakh

20%

Income above Rs 15 lakh

30%

Rs 12.5 lakh to Rs 15 lakh

25%

Income above Rs 15 lakh

30%

Under the New Tax Regime, many deductions and exemptions are not applicable. However, there are some exemptions and deductions available, such as:

1. Transport allowances for specially-abled individuals.
2. Conveyance allowance received for work-related travel expenses.
3. Compensation for travel costs during tours or transfers.
4. Daily allowance for regular expenses during absence from the regular place of work.

 

 

Exceptions to the Income Tax Slab:
Capital gains income is taxed differently based on the type of asset and how long it’s been owned.
Assets are categorized as long-term or short-term based on the holding period, which varies for different assets.
Here’s a summary of the holding period, asset type, and corresponding tax rates:

Financial Year (FY):
The financial year is the period from April 1st to March 31st during which income is earned and reported for taxation purposes. For example, FY 2022-23 starts from April 1st, 2022, to March 31st, 2023.

Assessment Year (AY):
The assessment year follows the financial year, spanning from April 1st to March 31st of the following year. Taxpayers assess their income earned during the financial year and pay taxes during this period. For instance, AY 2023-24 is for incomes earned in FY 2022-23.

Assessee:
An assessee is an individual or entity that assesses their income and pays taxes according to the Income Tax Act. This includes individuals, partnership firms, companies, Associations of Persons (AOPs), Trusts, etc.

 




What is PAN?
A PAN, or Permanent Account Number, is a special 10-digit alphanumeric code given by the Indian Income Tax Department to taxpayers. This number serves as a unique identifier for individuals and helps in tracking their tax-related activities. Whenever a person pays taxes or files returns, they must mention their PAN. Additionally, PAN is shared with banks, mutual funds, and other financial entities. This allows the Income Tax Department to monitor all financial transactions associated with that PAN. Essentially, PAN links an individual’s financial activities with the tax department, making it easier for tax authorities to manage taxation processes and track incomes.

What is TAN?
TAN, or Tax Deduction and Collection Account Number, is a unique 10-digit alphanumeric code issued by the Income Tax Department. It is primarily meant for entities responsible for deducting or collecting taxes at the source. Any entity involved in tax deduction (TDS) or tax collection (TCS) must obtain a TAN. This number is crucial for filing TDS/TCS returns, making tax payments, and issuing TDS/TCS certificates. By quoting TAN in all relevant documents and transactions, organizations ensure compliance with tax regulations and facilitate smooth processing of tax-related activities.

Residents and Non-residents
The levy of income tax in India hinges on the residential status of taxpayers. Residents are individuals who qualify under Indian tax laws and are liable to pay taxes on their worldwide income, including earnings from both within and outside India. Conversely, non-residents are only required to pay taxes on income earned within India. Determining residential status is crucial and must be done for each financial year separately. This classification ensures that tax liabilities are appropriately assessed based on the taxpayer’s presence and economic activities within the country’s jurisdiction.

Income Tax Payment
Tax Deducted at Source (TDS)
Tax deducted at source (TDS) is a mechanism where tax is deducted by the payer while making certain payments to the recipient. The recipient can later claim credit for this TDS amount by adjusting it against their final tax liability.

Advance Tax
Advance tax is tax paid in advance by taxpayers if their estimated tax liability for the year exceeds Rs 10,000. The government specifies due dates for paying advance tax installments throughout the financial year.

Self-Assessment Tax
Self-assessment tax is the balance tax that taxpayers have to pay on their assessed income after adjusting for advance tax and TDS. It is calculated based on the total income tax liability determined by the taxpayer.

E-Payment of Taxes
Taxpayers can conveniently pay advance tax and self-assessment tax online through the NSDL website. However, they need to have net banking facilities with authorized banks for e-payment.

Filing your ITR: A Simple Guide
E-filing your income tax return has become mandatory for most taxpayers, except for a few exceptions:

1. Taxpayers aged 80 and above are exempt from e-filing.
2. Taxpayers earning less than Rs 5 lakhs and not seeking a refund also do not need to e-file.
For everyone else, e-filing is compulsory. The deadline for filing returns is typically July 31 after the financial year ends.

If you miss the deadline, here are the consequences:

– You won’t be able to carry forward losses (except for house property loss) to future years.
– Your refund claims may be delayed.
– Obtaining home loans might become more challenging.
– Late filing fees of up to Rs 5,000 (for incomes above Rs 5 lakhs) or Rs 1,000 (for incomes below Rs 5 lakhs) may be levied under Section 234F.
– Interest may be charged under Section 234A if taxes are due as of July 31.

E-filing offers several benefits beyond just submitting your return. Platforms like Tax Craft help you maximize deductions and invest wisely. After filing online, you can either e-verify the return or print the ITR V and send it to CPC, Bengaluru, for processing.

Income Tax Return
Taxpayers must file their income tax return annually using the prescribed ITR forms provided by the income tax department. There are seven ITR forms available, and taxpayers need to select the appropriate form based on their income sources and other relevant factors.

 

Income Tax Forms: A Quick Overview
When filing your income tax return (ITR), you’ll need to choose the appropriate form based on your income sources and other factors. Here’s a simple breakdown of the seven ITR forms:

1. ITR-1: For individuals (residents) with income from salary, one house property, other sources, agricultural income less than Rs 5,000, and total income up to Rs 50 lakh.

2. ITR-2: For individuals/HUFs without any business or profession under proprietorship, but with more than one house property.

3. ITR-3: For individuals/HUFs with income from proprietary business or profession, and income as a partner in a firm.

4. ITR-4: For individuals/HUFs with presumptive income from business or profession, along with one house property.

5. ITR-5: For partnership firms or Limited Liability Partnerships (LLPs).

6. ITR-6: For companies.

7. ITR-7: For trusts.

Choose the form that best matches your income sources and follow the instructions carefully to complete your tax filing process.

 

Essential Documents for ITR Filing Simplified
Before filing your income tax return (ITR), make sure you have thes

    \n\n
    Context: \n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.8)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to handle user input
async def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)

    chain = await get_conversational_chain()

    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    st.write("Reply: ", response["output_text"])

def main():
    st.set_page_config("TAX CRAFT")

    st.title("TAX CRAFT")
    
    uploaded_files = st.file_uploader("Upload your Doc", accept_multiple_files=True, type=["pdf", "pptx"])
    user_question = st.text_input("Ask a Question from the Files")

    if user_question:
     asyncio.run(user_input(user_question))
    if st.button("Submit & Process"):
        with st.spinner("Processing..."):
            pdf_files = [file for file in uploaded_files if file.type == "application/pdf"]
            pptx_files = [file for file in uploaded_files if file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"]
                
            raw_text = get_pdf_text(pdf_files) + get_pptx_text(pptx_files)
            text_chunks = get_text_chunks(raw_text)
            get_vector_store(text_chunks)
            st.success("Done")

    with st.sidebar:
        
    
        
       
        st.write("<iframe src='https://www.chatbase.co/chatbot-iframe/pgFbsOUgVK4E4Xakpbj-L' width='100%' height='800px' style='border:none;'></iframe>", unsafe_allow_html=True)

    st.write("Created with ❤ by CODE RANGERS")

if _name_ == "_main_":
    main()
